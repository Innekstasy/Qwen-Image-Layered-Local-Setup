from diffusers import QwenImageLayeredPipeline
import torch
from PIL import Image
from pptx import Presentation
import os
import gradio as gr
import numpy as np
import random
import tempfile
import zipfile 
from psd_tools import PSDImage
MAX_SEED = np.iinfo(np.int32).max

MODEL_PATH = r"D:\AI\Qwen-Image-Layered\model"
ASSETS_PATH = r"D:\AI\Qwen-Image-Layered\repo\assets\test_images"

pipeline = QwenImageLayeredPipeline.from_pretrained(MODEL_PATH, torch_dtype=torch.bfloat16, device_map="balanced")
pipeline.set_progress_bar_config(disable=None)

def imagelist_to_pptx(img_files):
    with Image.open(img_files[0]) as img:
        img_width_px, img_height_px = img.size

    def px_to_emu(px, dpi=96):
        inch = px / dpi
        emu = inch * 914400
        return int(emu)

    prs = Presentation()
    prs.slide_width = px_to_emu(img_width_px)
    prs.slide_height = px_to_emu(img_height_px)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = top = 0
    for img_path in img_files:
        slide.shapes.add_picture(img_path, left, top, width=px_to_emu(img_width_px), height=px_to_emu(img_height_px))
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        prs.save(tmp.name)
        return tmp.name

def imagelist_to_psd(img_files):
    layers = []
    for path in img_files:
        layers.append(Image.open(path).convert('RGBA'))
    width, height = layers[0].size
    psd = PSDImage.new(mode='RGBA', size=(width, height))
    for i, img in enumerate(layers):
        name = f"Layer {i + 1}"
        layer = psd.create_pixel_layer(image=img, name=name)
        psd.append(layer)
    with tempfile.NamedTemporaryFile(suffix=".psd", delete=False) as tmp:
        psd.save(tmp.name)
        return tmp.name

def infer(input_image,
          seed=777,
          randomize_seed=False,
          prompt=None,
          neg_prompt=" ",
          true_guidance_scale=4.0,
          num_inference_steps=50,
          layer=4,
          cfg_norm=True,
          use_en_prompt=True):
    
    if randomize_seed:
        seed = random.randint(0, MAX_SEED)
    if isinstance(input_image, list):
        input_image = input_image[0]
    if isinstance(input_image, str):
        pil_image = Image.open(input_image).convert("RGB").convert("RGBA")
    elif isinstance(input_image, Image.Image):
        pil_image = input_image.convert("RGB").convert("RGBA")
    elif isinstance(input_image, np.ndarray):
        pil_image = Image.fromarray(input_image).convert("RGB").convert("RGBA")
    else:
        raise ValueError("Unsupported input_image type: %s" % type(input_image))
    
    inputs = {
        "image": pil_image,
        "generator": torch.Generator(device='cuda').manual_seed(seed),
        "true_cfg_scale": true_guidance_scale,
        "prompt": prompt,
        "negative_prompt": neg_prompt,
        "num_inference_steps": num_inference_steps,
        "num_images_per_prompt": 1,
        "layers": layer,
        "resolution": 640,
        "cfg_normalize": cfg_norm,
        "use_en_prompt": use_en_prompt,
    }
    print(inputs)
    with torch.inference_mode():
        output = pipeline(**inputs)
        output_images = output.images[0]
    
    output = []
    temp_files = []
    for i, image in enumerate(output_images):
        output.append(image)
        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        image.save(tmp.name)
        temp_files.append(tmp.name)
    
    pptx_path = imagelist_to_pptx(temp_files)
    with tempfile.NamedTemporaryFile(suffix=".zip", delete=False) as tmp:
        with zipfile.ZipFile(tmp.name, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for i, img_path in enumerate(temp_files):
                zipf.write(img_path, f"layer_{i+1}.png")
        zip_path = tmp.name
    psd_path = imagelist_to_psd(temp_files)
    return output, pptx_path, zip_path, psd_path

examples = [os.path.join(ASSETS_PATH, f"{i}.png") for i in range(1, 14)]

with gr.Blocks() as demo:
    with gr.Column(elem_id="col-container"):
        gr.HTML('''<p align="center"><img src="https://qianwen-res.oss-cn-beijing.aliyuncs.com/Qwen-Image/layered/qwen-image-layered-logo.png" width="400"/><p>''')
        gr.Markdown("""
                    The text prompt is intended to describe the overall content of the input image—including elements that may be partially occluded. It is not designed to control the semantic content of individual layers explicitly.
                    """)
        with gr.Row():
            with gr.Column(scale=1):
                input_image = gr.Image(label="Input Image", image_mode="RGBA")
                with gr.Accordion("Advanced Settings", open=False):
                    prompt = gr.Textbox(label="Prompt (Optional)", placeholder="Describe the image content (optional)", value="", lines=3)
                    neg_prompt = gr.Textbox(label="Negative Prompt (Optional)", value=" ", lines=3)
                    seed = gr.Slider(label="Seed", minimum=0, maximum=MAX_SEED, step=1, value=0)
                    randomize_seed = gr.Checkbox(label="Randomize seed", value=True)
                    true_guidance_scale = gr.Slider(label="True guidance scale", minimum=1.0, maximum=10.0, step=0.1, value=4.0)
                    num_inference_steps = gr.Slider(label="Number of inference steps", minimum=1, maximum=50, step=1, value=50)
                    layer = gr.Slider(label="Layers", minimum=2, maximum=10, step=1, value=4)
                    cfg_norm = gr.Checkbox(label="Enable CFG normalization", value=True)
                    use_en_prompt = gr.Checkbox(label="Auto-caption in English (uncheck for Chinese)", value=True)
                run_button = gr.Button("Decompose!", variant="primary")
            with gr.Column(scale=2):
                gallery = gr.Gallery(label="Layers", columns=4, rows=1, format="png")
                with gr.Row():
                    export_file = gr.File(label="Download PPTX")
                    export_zip_file = gr.File(label="Download ZIP")
                    export_psd_file = gr.File(label="Download PSD")

    gr.Examples(
        examples=examples,
        inputs=[input_image],
        outputs=[gallery, export_file, export_zip_file, export_psd_file],
        fn=infer,
        examples_per_page=14,
        cache_examples=False,
        run_on_click=True
    )

    run_button.click(
        fn=infer,
        inputs=[input_image, seed, randomize_seed, prompt, neg_prompt, true_guidance_scale, num_inference_steps, layer, cfg_norm, use_en_prompt],
        outputs=[gallery, export_file, export_zip_file, export_psd_file],
    )

demo.launch(
    server_name="0.0.0.0",
    server_port=7869,
)